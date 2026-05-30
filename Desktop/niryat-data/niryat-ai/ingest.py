"""
NiryatAI Document Ingestion Pipeline
Extracts text from PDF, DOCX, DOC, XLSX, XLT, PPTX, PPT, MP4, MP3.
Chunks, embeds with sentence-transformers, stores in ChromaDB.
"""

import os
import re
import subprocess
import hashlib
import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation
import chromadb
from chromadb.utils import embedding_functions

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(SCRIPT_DIR, "..")
CHROMA_PATH = os.path.join(SCRIPT_DIR, "backend", "chroma_db")
COLLECTION_NAME = "niryat_exim"
CHUNK_SIZE = 400
CHUNK_OVERLAP = 50

SKIP_FOLDERS = {"niryat-ai", ".git", "__pycache__", ".DS_Store"}


def extract_pdf(filepath: str) -> str:
    text = ""
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        print(f"  [PDF error] {e}")
    return text


def extract_xlsx(filepath: str) -> str:
    text = ""
    try:
        xls = pd.ExcelFile(filepath)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text += f"\n--- Sheet: {sheet_name} ---\n"
            text += df.to_string(index=False) + "\n"
    except Exception as e:
        print(f"  [XLSX error] {e}")
    return text


def extract_docx(filepath: str) -> str:
    text = ""
    try:
        doc = Document(filepath)
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text + "\n"
    except Exception as e:
        print(f"  [DOCX error] {e}")
    return text


def extract_doc(filepath: str) -> str:
    """Extract text from .doc files using textutil (macOS) or antiword."""
    text = ""
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", filepath],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0:
            text = result.stdout
        else:
            result2 = subprocess.run(
                ["antiword", filepath],
                capture_output=True, text=True, timeout=30,
            )
            if result2.returncode == 0:
                text = result2.stdout
    except FileNotFoundError:
        print("  [DOC] textutil/antiword not found, skipping .doc")
    except Exception as e:
        print(f"  [DOC error] {e}")
    return text


def extract_pptx(filepath: str) -> str:
    text = ""
    try:
        prs = Presentation(filepath)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_text.append(t)
            if slide_text:
                text += f"\n--- Slide {slide_num} ---\n"
                text += "\n".join(slide_text) + "\n"
    except Exception as e:
        print(f"  [PPTX error] {e}")
    return text


def extract_ppt(filepath: str) -> str:
    """Convert .ppt to text using LibreOffice or textutil on macOS."""
    text = ""
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", filepath],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            text = result.stdout
    except Exception as e:
        print(f"  [PPT error] {e}")
    return text


def transcribe_audio_video(filepath: str) -> str:
    """Transcribe MP4/MP3 using openai-whisper."""
    text = ""
    try:
        import whisper
        print(f"  Transcribing with Whisper (this may take a while)...")
        model = whisper.load_model("base")
        result = model.transcribe(filepath, language="en")
        text = result.get("text", "")
    except ImportError:
        print("  [Whisper] openai-whisper not installed, skipping audio/video")
    except Exception as e:
        print(f"  [Whisper error] {e}")
    return text


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
    ".xlt": extract_xlsx,
    ".docx": extract_docx,
    ".doc": extract_doc,
    ".pptx": extract_pptx,
    ".ppt": extract_ppt,
    ".mp4": transcribe_audio_video,
    ".mp3": transcribe_audio_video,
}


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    words = text.split()
    if len(words) <= chunk_size:
        return [text.strip()] if text.strip() else []
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks


def main():
    data_dir = os.path.abspath(DATA_DIR)
    print(f"Scanning: {data_dir}")

    if not os.path.exists(data_dir):
        print(f"Data directory not found: {data_dir}")
        return

    ef = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name="all-MiniLM-L6-v2"
    )
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        embedding_function=ef,
    )
    print(f"ChromaDB '{COLLECTION_NAME}' — existing docs: {collection.count()}")

    all_chunks = []
    all_ids = []
    all_metadatas = []
    seen_filenames = set()

    for root, dirs, files in os.walk(data_dir):
        dirs[:] = [d for d in dirs if d not in SKIP_FOLDERS]

        for filename in sorted(files):
            if filename.startswith("."):
                continue

            ext = os.path.splitext(filename)[1].lower()
            if ext not in EXTRACTORS:
                continue

            if filename in seen_filenames:
                print(f"Skipping duplicate: {filename}")
                continue
            seen_filenames.add(filename)

            filepath = os.path.join(root, filename)
            rel_folder = os.path.relpath(root, data_dir)
            print(f"Processing: {rel_folder}/{filename}")

            text = EXTRACTORS[ext](filepath)

            if not text.strip():
                print(f"  No text extracted, skipping.")
                continue

            text = re.sub(r"\n{3,}", "\n\n", text)
            chunks = chunk_text(text)
            print(f"  {len(text):,} chars → {len(chunks)} chunks")

            for i, chunk in enumerate(chunks):
                doc_id = hashlib.md5(f"{filename}_{i}".encode()).hexdigest()
                all_chunks.append(chunk)
                all_ids.append(doc_id)
                all_metadatas.append({
                    "source": filename,
                    "chunk_index": i,
                    "folder": rel_folder,
                })

    if not all_chunks:
        print("No documents found to ingest.")
        return

    print(f"\nIngesting {len(all_chunks)} chunks from {len(seen_filenames)} files...")

    batch_size = 100
    for i in range(0, len(all_chunks), batch_size):
        batch_end = min(i + batch_size, len(all_chunks))
        collection.upsert(
            ids=all_ids[i:batch_end],
            documents=all_chunks[i:batch_end],
            metadatas=all_metadatas[i:batch_end],
        )
        print(f"  Batch {i // batch_size + 1}: {batch_end}/{len(all_chunks)}")

    print(f"\nDone! Total docs in collection: {collection.count()}")


if __name__ == "__main__":
    main()
